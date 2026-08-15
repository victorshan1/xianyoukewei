#!/usr/bin/env python3
"""生成乡村课堂AI助教项目说明PPT - 精美升级版"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from lxml import etree
import os

# ── 精美配色方案 ──
# 主色：深蓝教育感
PRIMARY = RGBColor(0x0F, 0x4C, 0x81)       # 深蓝
PRIMARY_LIGHT = RGBColor(0x1E, 0x88, 0xE5) # 亮蓝
PRIMARY_DARK = RGBColor(0x0A, 0x33, 0x58)  # 深蓝黑
ACCENT = RGBColor(0x00, 0xB8, 0xA9)        # 青绿（教育活力）
ACCENT2 = RGBColor(0xFF, 0x98, 0x00)       # 暖橙（温暖关怀）
ACCENT3 = RGBColor(0x7C, 0x4D, 0xFF)       # 紫色（创新科技）
ACCENT4 = RGBColor(0x00, 0xC8, 0x53)       # 绿色（成长）
TEXT_DARK = RGBColor(0x1A, 0x1A, 0x2E)     # 深文字
TEXT_LIGHT = RGBColor(0x6B, 0x72, 0x80)    # 浅文字
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BG_GRAY = RGBColor(0xF5, 0xF7, 0xFA)       # 背景灰
BG_BLUE = RGBColor(0xE3, 0xF2, 0xFD)       # 浅蓝背景
BG_TEAL = RGBColor(0xE0, 0xF2, 0xF1)       # 浅青绿
BG_ORANGE = RGBColor(0xFF, 0xF3, 0xE0)     # 浅橙
BG_PURPLE = RGBColor(0xF3, 0xE8, 0xFF)     # 浅紫
BG_GREEN = RGBColor(0xE8, 0xF5, 0xE9)      # 浅绿

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
W = prs.slide_width
H = prs.slide_height

# ── 辅助函数 ──
def set_slide_bg(slide, color):
    """设置幻灯片背景色"""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_shape(slide, left, top, width, height, color, shape_type=MSO_SHAPE.RECTANGLE, transparency=0):
    """添加形状"""
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    if transparency > 0:
        # 设置透明度
        spPr = shape._element.spPr
        solidFill = spPr.find(qn('a:solidFill'))
        if solidFill is not None:
            srgbClr = solidFill.find(qn('a:srgbClr'))
            if srgbClr is not None:
                alpha = etree.SubElement(srgbClr, qn('a:alpha'))
                alpha.set('val', str(int((1 - transparency) * 100000)))
    return shape

def rgb_to_hex(rgb_color):
    """将RGBColor转换为十六进制字符串"""
    return str(rgb_color)

def add_gradient_shape(slide, left, top, width, height, color1, color2, angle=45, shape_type=MSO_SHAPE.RECTANGLE):
    """添加渐变形状"""
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    shape.line.fill.background()
    
    # 使用XML实现渐变
    spPr = shape._element.spPr
    # 清除现有填充
    for child in list(spPr):
        if child.tag.endswith('}solidFill') or child.tag.endswith('}gradFill') or child.tag.endswith('}noFill'):
            spPr.remove(child)
    
    gradFill = etree.SubElement(spPr, qn('a:gradFill'))
    gsLst = etree.SubElement(gradFill, qn('a:gsLst'))
    
    # 颜色1
    gs1 = etree.SubElement(gsLst, qn('a:gs'))
    gs1.set('pos', '0')
    srgb1 = etree.SubElement(gs1, qn('a:srgbClr'))
    srgb1.set('val', rgb_to_hex(color1))
    
    # 颜色2
    gs2 = etree.SubElement(gsLst, qn('a:gs'))
    gs2.set('pos', '100000')
    srgb2 = etree.SubElement(gs2, qn('a:srgbClr'))
    srgb2.set('val', rgb_to_hex(color2))
    
    lin = etree.SubElement(gradFill, qn('a:lin'))
    lin.set('ang', str(int(angle * 60000)))
    lin.set('scaled', '1')
    
    return shape

def add_text(slide, left, top, width, height, text, font_size=18, color=TEXT_DARK, 
             bold=False, alignment=PP_ALIGN.LEFT, font_name='Microsoft YaHei', anchor=MSO_ANCHOR.TOP):
    """添加文本框"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_multiline_text(slide, left, top, width, height, lines, font_size=14, color=TEXT_DARK, 
                       alignment=PP_ALIGN.LEFT, line_spacing=1.5, font_name='Microsoft YaHei'):
    """添加多行文本"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = font_name
        p.alignment = alignment
        p.space_after = Pt(font_size * (line_spacing - 1))
    return txBox

def add_card(slide, left, top, width, height, bg_color=WHITE, border_color=None, radius=0.08):
    """添加卡片（带圆角和阴影效果）"""
    # 卡片主体
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    card.fill.solid()
    card.fill.fore_color.rgb = bg_color
    if border_color:
        card.line.color.rgb = border_color
        card.line.width = Pt(1)
    else:
        card.line.fill.background()
    # 调整圆角
    card.adjustments[0] = radius
    return card

def add_top_accent(slide, color=PRIMARY_LIGHT):
    """顶部装饰条"""
    add_shape(slide, 0, 0, W, Inches(0.06), color)

def add_bottom_bar(slide, page_num, total=10):
    """底部信息栏"""
    add_shape(slide, 0, H - Inches(0.45), W, Inches(0.45), PRIMARY_DARK)
    add_text(slide, Inches(0.5), H - Inches(0.42), Inches(6), Inches(0.35),
             '乡村课堂AI助教  ·  小有可为AI向善创新挑战赛', 11, RGBColor(0xB0, 0xBE, 0xC5))
    add_text(slide, W - Inches(2.5), H - Inches(0.42), Inches(2), Inches(0.35),
             f'{page_num} / {total}', 11, RGBColor(0xB0, 0xBE, 0xC5), alignment=PP_ALIGN.RIGHT)

def add_page_title(slide, title, subtitle='', icon=''):
    """页面标题区域"""
    # 标题装饰线
    add_shape(slide, Inches(0.8), Inches(0.55), Inches(0.12), Inches(0.6), ACCENT)
    
    if icon:
        add_text(slide, Inches(1.1), Inches(0.45), Inches(0.7), Inches(0.7),
                 icon, 32, PRIMARY, bold=True)
        title_left = Inches(1.85)
    else:
        title_left = Inches(1.1)
    
    add_text(slide, title_left, Inches(0.42), Inches(10), Inches(0.6),
             title, 30, PRIMARY_DARK, bold=True)
    if subtitle:
        add_text(slide, title_left, Inches(1.0), Inches(10), Inches(0.4),
                 subtitle, 15, TEXT_LIGHT)

def add_decorative_circles(slide):
    """添加装饰性圆形元素"""
    # 右上角装饰
    c1 = add_shape(slide, Inches(12), Inches(-0.8), Inches(2.5), Inches(2.5), 
                   PRIMARY_LIGHT, MSO_SHAPE.OVAL, transparency=0.1)
    c2 = add_shape(slide, Inches(11.2), Inches(0.2), Inches(1.5), Inches(1.5), 
                   ACCENT, MSO_SHAPE.OVAL, transparency=0.15)

def add_stat_card(slide, left, top, number, label, color=PRIMARY_LIGHT, bg_color=WHITE):
    """数据统计卡片"""
    card = add_card(slide, left, top, Inches(2.3), Inches(1.3), bg_color)
    add_text(slide, left, top + Inches(0.15), Inches(2.3), Inches(0.6),
             number, 32, color, bold=True, alignment=PP_ALIGN.CENTER)
    add_text(slide, left, top + Inches(0.75), Inches(2.3), Inches(0.4),
             label, 13, TEXT_LIGHT, alignment=PP_ALIGN.CENTER)
    return card

# ═══════════════════════════════════════════════════════════
# 第1页：封面 - 大气渐变设计
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, WHITE)

# 左侧渐变区域
add_gradient_shape(slide, 0, 0, Inches(7.2), H, PRIMARY_DARK, PRIMARY_LIGHT, angle=135)

# 装饰圆形
add_shape(slide, Inches(-1.5), Inches(-1.5), Inches(5), Inches(5), 
          PRIMARY_LIGHT, MSO_SHAPE.OVAL, transparency=0.7)
add_shape(slide, Inches(4.5), Inches(5), Inches(4), Inches(4), 
          ACCENT, MSO_SHAPE.OVAL, transparency=0.8)
add_shape(slide, Inches(5.5), Inches(-0.5), Inches(2), Inches(2), 
          ACCENT2, MSO_SHAPE.OVAL, transparency=0.85)

# 左侧内容
add_text(slide, Inches(0.8), Inches(1.8), Inches(6), Inches(0.5),
         '🌱 小有可为 · AI向善创新挑战赛', 16, RGBColor(0x81, 0xD4, 0xFA))
add_text(slide, Inches(0.8), Inches(2.5), Inches(6), Inches(1.2),
         '乡村课堂AI助教', 52, WHITE, bold=True)
add_text(slide, Inches(0.8), Inches(3.8), Inches(6), Inches(0.6),
         '让每一个乡村孩子都能享受AI教育', 22, RGBColor(0xB3, 0xE5, 0xFC))

# 分隔线
add_shape(slide, Inches(0.8), Inches(4.7), Inches(2), Inches(0.06), ACCENT)

add_text(slide, Inches(0.8), Inches(5.0), Inches(6), Inches(0.4),
         '参赛赛道：乡村教育 · 大山里的AI课', 15, RGBColor(0x90, 0xCA, 0xF9))
add_text(slide, Inches(0.8), Inches(5.5), Inches(6), Inches(0.4),
         '作品类型：Web应用（纯前端架构）', 15, RGBColor(0x90, 0xCA, 0xF9))

# 右侧内容
add_text(slide, Inches(7.8), Inches(2.2), Inches(5), Inches(0.6),
         '三端协同  ·  AI赋能', 24, PRIMARY_DARK, bold=True)

# 右侧三个特色卡片
features_cover = [
    ('👩‍🏫 教师端', '智能备课 · 学生画像', PRIMARY),
    ('👨‍🎓 学生端', '拍照答疑 · 个性练习', ACCENT),
    ('👨‍👩‍👧 家长端', '学情报告 · 家校沟通', ACCENT2),
]
for i, (title, desc, color) in enumerate(features_cover):
    y = Inches(3.2 + i * 1.1)
    card = add_card(slide, Inches(7.8), y, Inches(4.8), Inches(0.9), BG_GRAY)
    # 左侧色条
    add_shape(slide, Inches(7.8), y, Inches(0.1), Inches(0.9), color)
    add_text(slide, Inches(8.1), y + Inches(0.1), Inches(4.3), Inches(0.4),
             title, 16, TEXT_DARK, bold=True)
    add_text(slide, Inches(8.1), y + Inches(0.5), Inches(4.3), Inches(0.35),
             desc, 13, TEXT_LIGHT)

# 底部标签
tags = ['纯前端', '多模态AI', '数据隔离', '低带宽友好', '多API兼容']
for i, tag in enumerate(tags):
    tag_shape = add_card(slide, Inches(7.8 + i * 1.05), Inches(6.2), Inches(0.95), Inches(0.4), 
                         BG_BLUE, radius=0.5)
    add_text(slide, Inches(7.8 + i * 1.05), Inches(6.25), Inches(0.95), Inches(0.3),
             tag, 10, PRIMARY, bold=True, alignment=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════
# 第2页：痛点分析 - 卡片式设计
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_GRAY)
add_top_accent(slide)
add_page_title(slide, '痛点分析', '乡村教育面临的真实困境', '📚')

pain_points = [
    ('📝', '教师备课负担重', '乡村教师身兼多科，缺乏AI辅助工具，教案质量参差不齐，备课时间长', ACCENT3, BG_PURPLE),
    ('🎯', '个性化教育缺失', '大班教学难以因材施教，学生薄弱知识点无法精准定位和针对性训练', ACCENT2, BG_ORANGE),
    ('💬', '家校沟通困难', '留守儿童多，家长外出务工，信息不对称，难以了解孩子学习情况', ACCENT, BG_TEAL),
    ('📶', '网络条件受限', '乡村网络带宽有限，传统云端应用难以流畅运行，使用门槛高', PRIMARY_LIGHT, BG_BLUE),
]

for i, (icon, title, desc, color, bg) in enumerate(pain_points):
    x = Inches(0.6 + i * 3.15)
    y = Inches(1.9)
    
    # 主卡片
    card = add_card(slide, x, y, Inches(2.9), Inches(4.8), WHITE)
    
    # 顶部图标区域
    add_shape(slide, x, y, Inches(2.9), Inches(1.4), color, MSO_SHAPE.ROUNDED_RECTANGLE)
    # 覆盖下半部分成方形
    add_shape(slide, x, y + Inches(0.7), Inches(2.9), Inches(0.7), color)
    
    add_text(slide, x, y + Inches(0.3), Inches(2.9), Inches(0.8),
             icon, 44, WHITE, alignment=PP_ALIGN.CENTER)
    
    # 标题
    add_text(slide, x + Inches(0.2), y + Inches(1.6), Inches(2.5), Inches(0.5),
             title, 19, TEXT_DARK, bold=True, alignment=PP_ALIGN.CENTER)
    
    # 分隔线
    add_shape(slide, x + Inches(1.1), y + Inches(2.2), Inches(0.7), Inches(0.04), color)
    
    # 描述
    add_text(slide, x + Inches(0.25), y + Inches(2.5), Inches(2.4), Inches(2),
             desc, 14, TEXT_LIGHT, alignment=PP_ALIGN.CENTER)

add_bottom_bar(slide, 2)

# ═══════════════════════════════════════════════════════════
# 第3页：产品简介 - 数据展示
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, WHITE)
add_top_accent(slide)
add_page_title(slide, '产品简介', '面向乡村教育的AI助教平台，三端口协同赋能', '🌟')

# 三个端口大卡片
ports = [
    ('👩\u200d🏫', '教师端', 'AI备课助手', '学生画像分析', '班级数据总览', '学生信息管理', PRIMARY, BG_BLUE),
    ('👨\u200d🎓', '学生端', '拍照智能答疑', '针对性练习', '自动错题本', '语音输入支持', ACCENT, BG_TEAL),
    ('👨\u200d👩\u200d👧', '家长端', '学情报告', '家校沟通话术', '绑定孩子数据', '权限安全隔离', ACCENT2, BG_ORANGE),
]

for i, (icon, title, *features, color, bg) in enumerate(ports):
    x = Inches(0.6 + i * 4.2)
    y = Inches(1.9)
    
    # 主卡片
    card = add_card(slide, x, y, Inches(3.9), Inches(3.8), WHITE)
    
    # 顶部渐变区域
    add_gradient_shape(slide, x, y, Inches(3.9), Inches(1.5), color, PRIMARY_LIGHT, 
                      angle=90, shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
    add_shape(slide, x, y + Inches(0.8), Inches(3.9), Inches(0.7), color)
    
    add_text(slide, x, y + Inches(0.25), Inches(3.9), Inches(0.6),
             icon, 40, WHITE, alignment=PP_ALIGN.CENTER)
    add_text(slide, x, y + Inches(0.95), Inches(3.9), Inches(0.45),
             title, 22, WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    
    # 功能列表
    for j, feature in enumerate(features):
        fy = y + Inches(1.8 + j * 0.48)
        add_text(slide, x + Inches(0.5), fy, Inches(0.3), Inches(0.35),
                 '✓', 16, color, bold=True)
        add_text(slide, x + Inches(0.9), fy, Inches(2.7), Inches(0.35),
                 feature, 14, TEXT_DARK)

# 数据统计区域
stats = [
    ('12,300+', '行代码', PRIMARY_LIGHT),
    ('7', '核心模块', ACCENT),
    ('3', '端口架构', ACCENT2),
    ('100%', '纯前端', ACCENT3),
]
for i, (num, label, color) in enumerate(stats):
    add_stat_card(slide, Inches(1 + i * 3), Inches(6.0), num, label, color)

add_bottom_bar(slide, 3)

# ═══════════════════════════════════════════════════════════
# 第4页：教师端功能
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_GRAY)
add_top_accent(slide)
add_page_title(slide, '核心功能 · 教师端', 'AI赋能教学全流程，减轻教师负担', '👩\u200d🏫')

features_teacher = [
    ('📝', 'AI备课助手', '智能生成完整教案\n自动分层作业设计\n基础/提高/拓展三层', PRIMARY),
    ('📊', '学生画像', '五维雷达图展示\n能力归因分析\n个性化教学建议', ACCENT3),
    ('📈', '班级总览', '成绩分布可视化\n学习趋势追踪\n薄弱知识点分析', ACCENT),
    ('👥', '学生管理', '批量导入导出\n完整信息管理\n学习档案记录', ACCENT2),
]

for i, (icon, title, desc, color) in enumerate(features_teacher):
    x = Inches(0.6 + i * 3.15)
    y = Inches(2.0)
    
    card = add_card(slide, x, y, Inches(2.9), Inches(4.5), WHITE)
    
    # 图标圆形背景
    icon_bg = add_shape(slide, x + Inches(0.85), y + Inches(0.4), Inches(1.2), Inches(1.2),
                        color, MSO_SHAPE.OVAL)
    add_text(slide, x + Inches(0.85), y + Inches(0.55), Inches(1.2), Inches(1),
             icon, 36, WHITE, alignment=PP_ALIGN.CENTER)
    
    add_text(slide, x + Inches(0.2), y + Inches(1.8), Inches(2.5), Inches(0.5),
             title, 19, TEXT_DARK, bold=True, alignment=PP_ALIGN.CENTER)
    
    # 分隔线
    add_shape(slide, x + Inches(1.05), y + Inches(2.4), Inches(0.8), Inches(0.04), color)
    
    add_text(slide, x + Inches(0.25), y + Inches(2.7), Inches(2.4), Inches(1.6),
             desc, 14, TEXT_LIGHT, alignment=PP_ALIGN.CENTER)

add_bottom_bar(slide, 4)

# ═══════════════════════════════════════════════════════════
# 第5页：学生端功能
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_GRAY)
add_top_accent(slide, ACCENT)
add_page_title(slide, '核心功能 · 学生端', 'AI陪伴式学习，随时答疑解惑', '👨\u200d🎓')

features_student = [
    ('📷', '拍照答疑', 
     '• 拍照/上传题目图片\n• AI智能识别与解析\n• 分步讲解解题思路\n• 支持语音输入问题\n• 通义千问VL多模态模型', ACCENT),
    ('🎯', '针对性练习', 
     '• 基于学生画像推荐\n• 薄弱点专项训练\n• 难度自适应调整\n• 知识点掌握追踪\n• 练习记录自动保存', PRIMARY_LIGHT),
    ('📕', '错题本', 
     '• 自动收录错题\n• 按知识点分类\n• 定期复习提醒\n• 错题重做功能\n• 掌握程度标记', ACCENT2),
]

for i, (icon, title, desc, color) in enumerate(features_student):
    x = Inches(0.6 + i * 4.2)
    y = Inches(2.0)
    
    card = add_card(slide, x, y, Inches(3.9), Inches(4.7), WHITE)
    
    # 顶部色条
    add_shape(slide, x, y, Inches(3.9), Inches(0.12), color, MSO_SHAPE.ROUNDED_RECTANGLE)
    add_shape(slide, x, y + Inches(0.06), Inches(3.9), Inches(0.06), color)
    
    # 图标
    icon_bg = add_shape(slide, x + Inches(1.35), y + Inches(0.4), Inches(1.2), Inches(1.2),
                        color, MSO_SHAPE.OVAL, transparency=0.15)
    add_text(slide, x + Inches(1.35), y + Inches(0.55), Inches(1.2), Inches(1),
             icon, 38, color, alignment=PP_ALIGN.CENTER)
    
    add_text(slide, x + Inches(0.3), y + Inches(1.8), Inches(3.3), Inches(0.5),
             title, 22, TEXT_DARK, bold=True, alignment=PP_ALIGN.CENTER)
    
    # 分隔线
    add_shape(slide, x + Inches(1.4), y + Inches(2.45), Inches(1.1), Inches(0.04), color)
    
    add_text(slide, x + Inches(0.4), y + Inches(2.75), Inches(3.1), Inches(1.8),
             desc, 14, TEXT_LIGHT)

add_bottom_bar(slide, 5)

# ═══════════════════════════════════════════════════════════
# 第6页：家长端功能
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_GRAY)
add_top_accent(slide, ACCENT2)
add_page_title(slide, '核心功能 · 家长端', '让关爱跨越距离，家校沟通无障碍', '👨\u200d👩\u200d👧')

features_parent = [
    ('📋', '学情报告', 
     '• 周报/月报自动生成\n• 可视化数据图表\n• 学习进度一目了然\n• 知识点掌握分析\n• 学习时长统计', ACCENT2),
    ('💬', '家校沟通', 
     '• AI生成沟通话术\n• 基于真实学习数据\n• 降低沟通门槛\n• 个性化沟通建议\n• 常用话术模板', ACCENT3),
]

for i, (icon, title, desc, color) in enumerate(features_parent):
    x = Inches(0.8 + i * 6.2)
    y = Inches(2.0)
    
    card = add_card(slide, x, y, Inches(5.7), Inches(3.8), WHITE)
    
    # 左侧图标区域
    add_shape(slide, x, y, Inches(1.8), Inches(3.8), color, MSO_SHAPE.ROUNDED_RECTANGLE)
    add_shape(slide, x + Inches(0.9), y, Inches(0.9), Inches(3.8), color)
    
    add_text(slide, x, y + Inches(1.2), Inches(1.8), Inches(1.2),
             icon, 56, WHITE, alignment=PP_ALIGN.CENTER)
    
    # 右侧内容
    add_text(slide, x + Inches(2.1), y + Inches(0.4), Inches(3.3), Inches(0.5),
             title, 24, TEXT_DARK, bold=True)
    add_shape(slide, x + Inches(2.1), y + Inches(1.0), Inches(1), Inches(0.04), color)
    add_text(slide, x + Inches(2.1), y + Inches(1.3), Inches(3.3), Inches(2.3),
             desc, 15, TEXT_LIGHT)

# 权限隔离说明卡片
security_card = add_card(slide, Inches(1.5), Inches(6.05), Inches(10.3), Inches(0.9), BG_BLUE)
add_shape(slide, Inches(1.5), Inches(6.05), Inches(0.12), Inches(0.9), PRIMARY_LIGHT)
add_text(slide, Inches(1.8), Inches(6.15), Inches(10), Inches(0.35),
         '🔒 权限安全隔离设计', 16, PRIMARY_DARK, bold=True)
add_text(slide, Inches(1.8), Inches(6.5), Inches(9.8), Inches(0.4),
         '家长端仅可查看已绑定孩子的学习数据，学生端仅可查看个人信息，教师端可查看全量数据，充分保障数据安全与隐私', 13, TEXT_LIGHT)

add_bottom_bar(slide, 6)

# ═══════════════════════════════════════════════════════════
# 第7页：技术亮点
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, WHITE)
add_top_accent(slide, ACCENT3)
add_page_title(slide, '技术亮点', '创新技术方案，专为乡村教育场景优化', '⚡')

highlights = [
    ('🤖', '多模态视觉AI', 
     '采用通义千问VL Max视觉理解模型\n支持图片OCR、数学公式识别\n分步解题推理，准确率高', PRIMARY),
    ('🔄', '多API兼容架构', 
     '支持阿里云通义千问、DeepSeek\n及任意OpenAI兼容API\n用户可自由切换和自定义模型', ACCENT),
    ('🔐', '三端权限隔离', 
     '教师端：全量数据权限\n学生端：仅个人信息\n家长端：仅绑定孩子数据', ACCENT3),
    ('📶', '低带宽优化', 
     '纯前端架构，无需后端服务器\n数据存储在浏览器本地\n弱网环境也能流畅使用', ACCENT2),
]

for i, (icon, title, desc, color) in enumerate(highlights):
    x = Inches(0.6 + i * 3.15)
    y = Inches(1.9)
    
    # 卡片带悬浮效果
    card = add_card(slide, x, y, Inches(2.9), Inches(4.8), WHITE)
    
    # 顶部渐变条
    add_gradient_shape(slide, x, y, Inches(2.9), Inches(0.15), color, PRIMARY_LIGHT, angle=0)
    
    # 图标区域
    icon_bg = add_shape(slide, x + Inches(0.95), y + Inches(0.5), Inches(1), Inches(1),
                        color, MSO_SHAPE.OVAL)
    add_text(slide, x + Inches(0.95), y + Inches(0.6), Inches(1), Inches(0.8),
             icon, 32, WHITE, alignment=PP_ALIGN.CENTER)
    
    add_text(slide, x + Inches(0.15), y + Inches(1.7), Inches(2.6), Inches(0.5),
             title, 18, TEXT_DARK, bold=True, alignment=PP_ALIGN.CENTER)
    
    add_shape(slide, x + Inches(0.95), y + Inches(2.3), Inches(1), Inches(0.04), color)
    
    add_text(slide, x + Inches(0.2), y + Inches(2.6), Inches(2.5), Inches(2),
             desc, 13, TEXT_LIGHT, alignment=PP_ALIGN.CENTER)

add_bottom_bar(slide, 7)

# ═══════════════════════════════════════════════════════════
# 第8页：技术架构
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_GRAY)
add_top_accent(slide)
add_page_title(slide, '技术架构', '纯前端架构，开箱即用，零部署成本', '🏗️')

# 架构图
# 用户层
user_layer = add_card(slide, Inches(1.5), Inches(2.0), Inches(10.3), Inches(0.9), PRIMARY)
add_text(slide, Inches(1.5), Inches(2.1), Inches(10.3), Inches(0.7),
         '👩\u200d🏫 教师端    👨\u200d🎓 学生端    👨\u200d👩\u200d👧 家长端', 
         20, WHITE, bold=True, alignment=PP_ALIGN.CENTER)

# 箭头
for i in range(3):
    ax = Inches(5 + i * 1.2)
    add_text(slide, ax, Inches(3.0), Inches(0.8), Inches(0.4),
             '▼', 20, TEXT_LIGHT, alignment=PP_ALIGN.CENTER)

# 应用层
app_card = add_card(slide, Inches(1.5), Inches(3.5), Inches(10.3), Inches(1.3), WHITE)
add_text(slide, Inches(1.8), Inches(3.6), Inches(2), Inches(0.35),
         '应用层', 15, PRIMARY, bold=True)
add_shape(slide, Inches(1.8), Inches(3.95), Inches(0.8), Inches(0.04), ACCENT)

app_modules = ['路由管理', '数据存储', 'API服务', '权限控制', 'UI组件库']
for i, mod in enumerate(app_modules):
    mx = Inches(1.8 + i * 1.95)
    mod_card = add_card(slide, mx, Inches(4.15), Inches(1.7), Inches(0.5), BG_BLUE, radius=0.3)
    add_text(slide, mx, Inches(4.22), Inches(1.7), Inches(0.4),
             mod, 13, PRIMARY, bold=True, alignment=PP_ALIGN.CENTER)

# 箭头
for i in range(3):
    ax = Inches(5 + i * 1.2)
    add_text(slide, ax, Inches(4.9), Inches(0.8), Inches(0.4),
             '▼', 20, TEXT_LIGHT, alignment=PP_ALIGN.CENTER)

# 技术栈层
tech_card = add_card(slide, Inches(1.5), Inches(5.4), Inches(10.3), Inches(1.5), WHITE)
add_text(slide, Inches(1.8), Inches(5.5), Inches(2), Inches(0.35),
         '技术栈', 15, PRIMARY, bold=True)
add_shape(slide, Inches(1.8), Inches(5.85), Inches(0.8), Inches(0.04), ACCENT)

tech_items = [
    'HTML5 + CSS3 + ES6+',
    'IndexedDB 本地数据库',
    'localStorage 配置存储',
    'OpenAI 兼容 API',
    'Canvas 数据可视化',
    'Web Speech API',
    'FileReader API',
    '响应式设计',
]
for i, item in enumerate(tech_items):
    col = i % 4
    row = i // 4
    tx = Inches(1.8 + col * 2.45)
    ty = Inches(6.0 + row * 0.4)
    add_text(slide, tx, ty, Inches(2.3), Inches(0.35),
             f'✦ {item}', 12, TEXT_DARK)

add_bottom_bar(slide, 8)

# ═══════════════════════════════════════════════════════════
# 第9页：使用场景
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, WHITE)
add_top_accent(slide, ACCENT4)
add_page_title(slide, '使用场景', '真实乡村课堂应用案例', '🎬')

scenarios = [
    ('🌅', '场景一：高效备课', 
     '乡村教师王老师要准备明天的数学课\n打开AI备课助手，输入课题\nAI自动生成完整教案和三层分层作业\n原来2小时的备课现在只需10分钟',
     PRIMARY, BG_BLUE),
    ('📱', '场景二：随时答疑', 
     '学生小明遇到一道几何题不会做\n用手机拍照上传题目\nAI识别图形并分步讲解解题思路\n就像有一个24小时在线的辅导老师',
     ACCENT, BG_TEAL),
    ('💝', '场景三：家校连心', 
     '在外打工的李妈妈想了解孩子学习\n打开家长端查看学情报告\nAI还帮她生成与老师沟通的话术\n让关爱不再因距离而缺席',
     ACCENT2, BG_ORANGE),
]

for i, (icon, title, desc, color, bg) in enumerate(scenarios):
    x = Inches(0.6 + i * 4.2)
    y = Inches(1.9)
    
    card = add_card(slide, x, y, Inches(3.9), Inches(4.9), WHITE)
    
    # 顶部图标背景
    add_shape(slide, x, y, Inches(3.9), Inches(1.6), color, MSO_SHAPE.ROUNDED_RECTANGLE)
    add_shape(slide, x, y + Inches(0.8), Inches(3.9), Inches(0.8), color)
    
    add_text(slide, x, y + Inches(0.3), Inches(3.9), Inches(0.8),
             icon, 44, WHITE, alignment=PP_ALIGN.CENTER)
    add_text(slide, x, y + Inches(1.1), Inches(3.9), Inches(0.4),
             title, 18, WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    
    # 描述内容
    add_text(slide, x + Inches(0.35), y + Inches(2.0), Inches(3.2), Inches(2.6),
             desc, 14, TEXT_LIGHT, alignment=PP_ALIGN.CENTER)

add_bottom_bar(slide, 9)

# ═══════════════════════════════════════════════════════════
# 第10页：成果与展望
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, WHITE)

# 顶部渐变区域
add_gradient_shape(slide, 0, 0, W, Inches(2.2), PRIMARY_DARK, PRIMARY_LIGHT, angle=135)
add_shape(slide, 0, Inches(1.1), W, Inches(1.1), PRIMARY_LIGHT)

add_text(slide, Inches(0.8), Inches(0.5), Inches(10), Inches(0.5),
         '🎯', 32, WHITE)
add_text(slide, Inches(0.8), Inches(1.0), Inches(11), Inches(0.8),
         '项目成果与未来规划', 34, WHITE, bold=True)
add_text(slide, Inches(0.8), Inches(1.75), Inches(11), Inches(0.4),
         '脚踏实地，展望未来，用AI点亮乡村教育', 16, RGBColor(0xB3, 0xE5, 0xFC))

# 左侧：项目成果
left_card = add_card(slide, Inches(0.6), Inches(2.6), Inches(5.9), Inches(3.9), BG_BLUE)
add_shape(slide, Inches(0.6), Inches(2.6), Inches(5.9), Inches(0.1), PRIMARY_LIGHT)

add_text(slide, Inches(0.9), Inches(2.8), Inches(5), Inches(0.5),
         '✅  项目成果', 22, PRIMARY_DARK, bold=True)

achievements = [
    '12,300+ 行代码，功能完整',
    '三端口架构：教师/学生/家长',
    '7大核心功能模块',
    '支持阿里云/DeepSeek/自定义API',
    '三端权限隔离，数据安全合规',
    '纯前端实现，开箱即用零部署',
    '适配桌面端和移动端',
    '低带宽友好，适合乡村环境',
]
for i, item in enumerate(achievements):
    add_text(slide, Inches(1.1), Inches(3.5 + i * 0.36), Inches(5.2), Inches(0.35),
             f'✓  {item}', 14, TEXT_DARK)

# 右侧：未来规划
right_card = add_card(slide, Inches(6.8), Inches(2.6), Inches(5.9), Inches(3.9), BG_GREEN)
add_shape(slide, Inches(6.8), Inches(2.6), Inches(5.9), Inches(0.1), ACCENT4)

add_text(slide, Inches(7.1), Inches(2.8), Inches(5), Inches(0.5),
         '🚀  未来规划', 22, RGBColor(0x1B, 0x5E, 0x20), bold=True)

future = [
    '接入更多学科题库资源',
    '支持完整离线模式，无网也能用',
    '增加教师间教案共享社区',
    '引入AI口语练习功能',
    '对接学校教务管理系统',
    '多语言支持（少数民族地区）',
    '家长端增加学习提醒推送',
    '开源社区共建，持续迭代',
]
for i, item in enumerate(future):
    add_text(slide, Inches(7.3), Inches(3.5 + i * 0.36), Inches(5.2), Inches(0.35),
             f'→  {item}', 14, TEXT_DARK)

# 底部标语
add_text(slide, Inches(1), Inches(6.7), Inches(11.3), Inches(0.5),
         '用AI点亮乡村教育，让每个孩子都能看见未来 ✨', 24, PRIMARY, bold=True, alignment=PP_ALIGN.CENTER)

add_bottom_bar(slide, 10)

# ── 保存 ──
output_path = '/workspace/乡村课堂AI助教-项目说明.pptx'
prs.save(output_path)
print(f'PPT 已生成：{output_path}')
print(f'文件大小：{os.path.getsize(output_path) / 1024:.1f} KB')
print(f'共 {len(prs.slides)} 页')
