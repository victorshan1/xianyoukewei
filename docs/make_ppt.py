#!/usr/bin/env python3
"""生成乡村课堂AI助教比赛PPT"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# === 颜色主题 ===
ORANGE_PRIMARY = RGBColor(0xFF, 0x6B, 0x35)       # 暖橙主色
ORANGE_LIGHT = RGBColor(0xFF, 0x8C, 0x42)          # 暖橙浅色
ORANGE_BG = RGBColor(0xFF, 0xF5, 0xF0)             # 暖橙背景
DARK_TEXT = RGBColor(0x2D, 0x2D, 0x2D)             # 深色文字
GRAY_TEXT = RGBColor(0x66, 0x66, 0x66)             # 灰色文字
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_BG = RGBColor(0xFA, 0xF9, 0xF7)              # 浅米色背景
GREEN = RGBColor(0x4E, 0xCD, 0xC4)                 # 青绿
PURPLE = RGBColor(0xA7, 0x8B, 0xFA)                # 淡紫
RED_ACCENT = RGBColor(0xFF, 0x4D, 0x4D)            # 红色强调

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
W = prs.slide_width
H = prs.slide_height


def add_bg(slide, color=LIGHT_BG):
    """添加纯色背景"""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape_bg(slide, left, top, width, height, color, alpha=None):
    """添加矩形色块"""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    if alpha is not None:
        from lxml import etree
        solidFill = shape.fill._fill
        srgbClr = solidFill.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}srgbClr')
        if srgbClr is not None:
            alpha_elem = etree.SubElement(srgbClr, '{http://schemas.openxmlformats.org/drawingml/2006/main}alpha')
            alpha_elem.set('val', str(int(alpha * 1000)))
    return shape


def add_rounded_rect(slide, left, top, width, height, color):
    """添加圆角矩形"""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_text_box(slide, left, top, width, height, text, font_size=18, color=DARK_TEXT, bold=False, alignment=PP_ALIGN.LEFT, font_name='Microsoft YaHei'):
    """添加文本框"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_multi_text(slide, left, top, width, height, lines, default_size=16, default_color=DARK_TEXT):
    """添加多行文本"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line_info in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        if isinstance(line_info, dict):
            p.text = line_info.get('text', '')
            p.font.size = Pt(line_info.get('size', default_size))
            p.font.color.rgb = line_info.get('color', default_color)
            p.font.bold = line_info.get('bold', False)
            p.font.name = line_info.get('font', 'Microsoft YaHei')
            p.alignment = line_info.get('align', PP_ALIGN.LEFT)
            p.space_before = Pt(line_info.get('space_before', 6))
        else:
            p.text = str(line_info)
            p.font.size = Pt(default_size)
            p.font.color.rgb = default_color
            p.font.name = 'Microsoft YaHei'
            p.space_before = Pt(6)
    return txBox


def add_icon_card(slide, left, top, width, height, icon, title, desc, bg_color=WHITE):
    """添加图标卡片"""
    card = add_rounded_rect(slide, left, top, width, height, bg_color)
    # 阴影效果通过边框模拟
    card.line.color.rgb = RGBColor(0xEE, 0xEE, 0xEE)
    card.line.width = Pt(1)

    # 图标
    add_text_box(slide, left + Inches(0.3), top + Inches(0.2), width - Inches(0.6), Inches(0.6),
                 icon, font_size=28, alignment=PP_ALIGN.CENTER)
    # 标题
    add_text_box(slide, left + Inches(0.2), top + Inches(0.75), width - Inches(0.4), Inches(0.4),
                 title, font_size=14, bold=True, color=ORANGE_PRIMARY, alignment=PP_ALIGN.CENTER)
    # 描述
    add_text_box(slide, left + Inches(0.15), top + Inches(1.15), width - Inches(0.3), height - Inches(1.3),
                 desc, font_size=11, color=GRAY_TEXT, alignment=PP_ALIGN.CENTER)
    return card


def add_top_bar(slide, title_text):
    """添加顶部装饰条"""
    add_shape_bg(slide, 0, 0, W, Inches(0.06), ORANGE_PRIMARY)
    # 标题
    add_text_box(slide, Inches(0.8), Inches(0.25), Inches(10), Inches(0.6),
                 title_text, font_size=22, bold=True, color=DARK_TEXT)
    # 装饰线
    add_shape_bg(slide, Inches(0.8), Inches(0.85), Inches(1.5), Inches(0.04), ORANGE_PRIMARY)


def add_page_number(slide, num, total=11):
    """添加页码"""
    add_text_box(slide, W - Inches(1.5), H - Inches(0.5), Inches(1), Inches(0.3),
                 f'{num}/{total}', font_size=10, color=GRAY_TEXT, alignment=PP_ALIGN.RIGHT)


# ============================================================
# 第1页：封面
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_bg(slide, WHITE)

# 顶部装饰线
add_shape_bg(slide, 0, 0, W, Inches(0.08), ORANGE_PRIMARY)

# 左侧大标题区域
add_text_box(slide, Inches(1.2), Inches(1.5), Inches(7), Inches(1.2),
             '乡村课堂AI助教', font_size=48, bold=True, color=DARK_TEXT)

add_text_box(slide, Inches(1.2), Inches(2.7), Inches(7), Inches(0.8),
             '让每一位乡村教师都有AI助手', font_size=24, color=ORANGE_PRIMARY)

# 分隔线
add_shape_bg(slide, Inches(1.2), Inches(3.6), Inches(2), Inches(0.04), ORANGE_PRIMARY)

# 副标题信息
add_multi_text(slide, Inches(1.2), Inches(3.9), Inches(6), Inches(2), [
    {'text': '阿里巴巴"小有可为"比赛 — 乡村教育赛道', 'size': 16, 'color': GRAY_TEXT, 'space_before': 8},
    {'text': '', 'size': 8},
    {'text': '纯前端 · 离线可用 · 三端协同 · 零部署成本', 'size': 14, 'color': ORANGE_LIGHT, 'bold': True, 'space_before': 12},
])

# 右侧装饰区域 - 模拟界面卡片
card = add_rounded_rect(slide, Inches(8.5), Inches(1.2), Inches(4.2), Inches(5.5), ORANGE_BG)
card.line.color.rgb = RGBColor(0xFF, 0xE0, 0xD0)
card.line.width = Pt(1)

# 卡片内模拟界面元素
add_text_box(slide, Inches(8.8), Inches(1.5), Inches(3.6), Inches(0.4),
             '🏫 备课助手 - 乡村课堂AI助教', font_size=12, bold=True, color=DARK_TEXT)
add_shape_bg(slide, Inches(8.8), Inches(2.0), Inches(3.6), Inches(0.03), RGBColor(0xDD, 0xDD, 0xDD))

# 模拟功能卡片
for i, (icon, name) in enumerate([('📚', '智能备课'), ('📊', '学生画像'), ('🏫', '班级总览')]):
    y = Inches(2.3) + Inches(i * 1.1)
    mini_card = add_rounded_rect(slide, Inches(8.8), y, Inches(3.6), Inches(0.9), WHITE)
    mini_card.line.color.rgb = RGBColor(0xEE, 0xEE, 0xEE)
    add_text_box(slide, Inches(8.9), y + Inches(0.1), Inches(0.5), Inches(0.5),
                 icon, font_size=20)
    add_text_box(slide, Inches(9.5), y + Inches(0.2), Inches(2.5), Inches(0.4),
                 name, font_size=13, bold=True, color=DARK_TEXT)

# 底部信息
add_text_box(slide, Inches(1.2), H - Inches(0.8), Inches(5), Inches(0.4),
             '技术有温度 · 教育无距离', font_size=12, color=GRAY_TEXT)
add_page_number(slide, 1)


# ============================================================
# 第2页：问题 — 乡村教育的真实困境
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_top_bar(slide, '乡村教育的真实困境')
add_page_number(slide, 2)

# 数据卡片
data_items = [
    ('1:40+', '乡村师生比', '远超城市 1:20', RGBColor(0xFF, 0xE8, 0xE0)),
    ('12h', '教师周备课时间', '是城市的 1.5 倍', RGBColor(0xE0, 0xF0, 0xFF)),
    ('67%', '学校无心理老师', '专业资源严重不足', RGBColor(0xE8, 0xFF, 0xE0)),
    ('40%+', '留守儿童占比', '家庭教育严重缺位', RGBColor(0xF0, 0xE8, 0xFF)),
]

for i, (num, label, desc, bg) in enumerate(data_items):
    x = Inches(0.8) + Inches(i * 3.1)
    card = add_rounded_rect(slide, x, Inches(1.2), Inches(2.8), Inches(2.2), bg)
    card.line.fill.background()
    add_text_box(slide, x + Inches(0.2), Inches(1.4), Inches(2.4), Inches(0.8),
                 num, font_size=36, bold=True, color=ORANGE_PRIMARY, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x + Inches(0.2), Inches(2.2), Inches(2.4), Inches(0.4),
                 label, font_size=14, bold=True, color=DARK_TEXT, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x + Inches(0.2), Inches(2.7), Inches(2.4), Inches(0.4),
                 desc, font_size=11, color=GRAY_TEXT, alignment=PP_ALIGN.CENTER)

# 三个痛点
pain_points = [
    ('👩‍🏫', '教师', '备课负担重\n因材施教难\n一人多科多年级'),
    ('', '学生', '问题无人及时解答\n错题无人整理\n学习信心不足'),
    ('👵', '家长', '想了解孩子学习\n但看不懂报告\n沟通渠道有限'),
]

for i, (icon, title, desc) in enumerate(pain_points):
    x = Inches(0.8) + Inches(i * 4.1)
    card = add_rounded_rect(slide, x, Inches(3.8), Inches(3.8), Inches(3.2), WHITE)
    card.line.color.rgb = RGBColor(0xEE, 0xEE, 0xEE)
    card.line.width = Pt(1)

    add_text_box(slide, x + Inches(0.3), Inches(4.0), Inches(0.6), Inches(0.6),
                 icon, font_size=32)
    add_text_box(slide, x + Inches(1.0), Inches(4.05), Inches(2.5), Inches(0.4),
                 title, font_size=18, bold=True, color=ORANGE_PRIMARY)
    add_shape_bg(slide, x + Inches(0.3), Inches(4.6), Inches(3.2), Inches(0.03), RGBColor(0xFF, 0xE8, 0xE0))

    for j, line in enumerate(desc.split('\n')):
        add_text_box(slide, x + Inches(0.3), Inches(4.8) + Inches(j * 0.45), Inches(3.2), Inches(0.4),
                     f'• {line}', font_size=13, color=DARK_TEXT)


# ============================================================
# 第3页：我们的方案
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_top_bar(slide, '我们的方案')
add_page_number(slide, 3)

# 一句话定位
add_text_box(slide, Inches(0.8), Inches(1.2), Inches(11), Inches(0.6),
             '一个纯前端的 AI 助教系统，教师、学生、家长三端协同', font_size=20, bold=True, color=DARK_TEXT)

add_text_box(slide, Inches(0.8), Inches(1.8), Inches(11), Inches(0.5),
             '核心定位：不是替代老师，而是让老师从重复劳动中解放出来', font_size=14, color=ORANGE_PRIMARY)

# 三个端口卡片
ports = [
    ('👩‍🏫', '教师端', ORANGE_PRIMARY, RGBColor(0xFF, 0xF5, 0xF0),
     ['智能备课（4种模板）', '学生画像（五维雷达图）', '班级总览（知识点热力图）', '成长轨迹时间轴']),
    ('👦', '学生端', GREEN, RGBColor(0xF0, 0xFD, 0xFB),
     ['拍照答疑（分步解析）', '针对性练习', '错题本（自动归类）', '知识点卡片']),
    ('👨‍👩‍👧', '家长端', PURPLE, RGBColor(0xF8, 0xF6, 0xFF),
     ['学情报告（周/月）', '教师建议', '成长追踪', '通俗易懂的可视化']),
]

for i, (icon, title, color, bg, features) in enumerate(ports):
    x = Inches(0.8) + Inches(i * 4.1)
    card = add_rounded_rect(slide, x, Inches(2.6), Inches(3.8), Inches(4.5), bg)
    card.line.color.rgb = color
    card.line.width = Pt(1.5)

    # 顶部色条
    add_shape_bg(slide, x + Inches(0.05), Inches(2.65), Inches(3.7), Inches(0.06), color)

    add_text_box(slide, x + Inches(0.3), Inches(2.9), Inches(0.6), Inches(0.6),
                 icon, font_size=36)
    add_text_box(slide, x + Inches(1.0), Inches(2.95), Inches(2.5), Inches(0.5),
                 title, font_size=22, bold=True, color=color)

    add_shape_bg(slide, x + Inches(0.3), Inches(3.6), Inches(3.2), Inches(0.03), color)

    for j, feat in enumerate(features):
        add_text_box(slide, x + Inches(0.3), Inches(3.8) + Inches(j * 0.6), Inches(3.2), Inches(0.5),
                     f'✓ {feat}', font_size=13, color=DARK_TEXT)


# ============================================================
# 第4页：技术架构
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_top_bar(slide, '技术架构')
add_page_number(slide, 4)

# 架构图 - 三层
layers = [
    ('🖥️ 展示层', 'HTML / CSS / JavaScript\nECharts 数据可视化\nFont Awesome 图标', ORANGE_PRIMARY),
    ('️ 逻辑层', '路由管理 / 状态管理\nIndexedDB 本地存储\nService Worker 离线缓存', GREEN),
    ('🤖 AI 层', 'Qwen API 智能生成\n自然语言理解\n多轮对话能力', PURPLE),
]

for i, (title, desc, color) in enumerate(layers):
    y = Inches(1.3) + Inches(i * 1.8)
    card = add_rounded_rect(slide, Inches(1.5), y, Inches(5), Inches(1.5), WHITE)
    card.line.color.rgb = color
    card.line.width = Pt(2)

    add_text_box(slide, Inches(1.8), y + Inches(0.15), Inches(4.5), Inches(0.4),
                 title, font_size=16, bold=True, color=color)
    add_text_box(slide, Inches(1.8), y + Inches(0.6), Inches(4.5), Inches(0.8),
                 desc, font_size=12, color=DARK_TEXT)

    # 箭头
    if i < 2:
        add_text_box(slide, Inches(3.8), y + Inches(1.5), Inches(0.5), Inches(0.4),
                     '▼', font_size=16, color=GRAY_TEXT, alignment=PP_ALIGN.CENTER)

# 右侧特性列表
features = [
    ('', '部署零门槛', '一台普通电脑即可运行\n无需服务器、无需数据库'),
    ('📡', '离线可用', 'Service Worker 缓存\n断网也能访问核心功能'),
    ('🔒', '数据安全', '所有数据存储在本地\n不上传云端，保护隐私'),
    ('📦', '纯前端架构', 'HTML/CSS/JS + IndexedDB\n零后端依赖'),
]

for i, (icon, title, desc) in enumerate(features):
    y = Inches(1.3) + Inches(i * 1.5)
    card = add_rounded_rect(slide, Inches(7.5), y, Inches(5.2), Inches(1.3), ORANGE_BG)
    card.line.fill.background()

    add_text_box(slide, Inches(7.7), y + Inches(0.1), Inches(0.5), Inches(0.5),
                 icon, font_size=22)
    add_text_box(slide, Inches(8.3), y + Inches(0.1), Inches(4), Inches(0.35),
                 title, font_size=14, bold=True, color=ORANGE_PRIMARY)
    add_text_box(slide, Inches(8.3), y + Inches(0.5), Inches(4), Inches(0.7),
                 desc, font_size=11, color=GRAY_TEXT)


# ============================================================
# 第5页：核心功能演示
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_top_bar(slide, '核心功能演示')
add_page_number(slide, 5)

# 教师端
add_text_box(slide, Inches(0.8), Inches(1.2), Inches(3), Inches(0.4),
             '👩‍🏫 教师端', font_size=18, bold=True, color=ORANGE_PRIMARY)
add_shape_bg(slide, Inches(0.8), Inches(1.65), Inches(1), Inches(0.04), ORANGE_PRIMARY)

teacher_feats = [
    ' 4种教案模板一键生成',
    '✏️ 在线富文本编辑',
    '📊 五维雷达图 + 时间轴',
    '🔥 班级知识点热力图',
]
for i, f in enumerate(teacher_feats):
    add_text_box(slide, Inches(0.8), Inches(1.9) + Inches(i * 0.55), Inches(5), Inches(0.45),
                 f, font_size=13, color=DARK_TEXT)

# 学生端
add_text_box(slide, Inches(4.8), Inches(1.2), Inches(3), Inches(0.4),
             ' 学生端', font_size=18, bold=True, color=GREEN)
add_shape_bg(slide, Inches(4.8), Inches(1.65), Inches(1), Inches(0.04), GREEN)

student_feats = [
    '📷 拍照上传，AI分步解析',
    '📇 知识点卡片关联',
    '📝 错题本自动归类',
    '✅ 标记掌握状态',
]
for i, f in enumerate(student_feats):
    add_text_box(slide, Inches(4.8), Inches(1.9) + Inches(i * 0.55), Inches(5), Inches(0.45),
                 f, font_size=13, color=DARK_TEXT)

# 家长端
add_text_box(slide, Inches(8.8), Inches(1.2), Inches(3), Inches(0.4),
             '‍👩‍👧 家长端', font_size=18, bold=True, color=PURPLE)
add_shape_bg(slide, Inches(8.8), Inches(1.65), Inches(1), Inches(0.04), PURPLE)

parent_feats = [
    '📋 周/月学情报告',
    '💡 教师建议通俗易懂',
    '📈 成绩趋势可视化',
    '🌟 成长追踪记录',
]
for i, f in enumerate(parent_feats):
    add_text_box(slide, Inches(8.8), Inches(1.9) + Inches(i * 0.55), Inches(5), Inches(0.45),
                 f, font_size=13, color=DARK_TEXT)

# 底部截图提示区域
screenshot_area = add_rounded_rect(slide, Inches(0.8), Inches(4.5), Inches(11.7), Inches(2.7), RGBColor(0xF5, 0xF5, 0xF5))
screenshot_area.line.color.rgb = RGBColor(0xDD, 0xDD, 0xDD)
screenshot_area.line.width = Pt(1)

add_text_box(slide, Inches(1.2), Inches(4.7), Inches(11), Inches(0.5),
             '📸 此处插入功能截图或动图', font_size=16, bold=True, color=GRAY_TEXT, alignment=PP_ALIGN.CENTER)

add_multi_text(slide, Inches(1.2), Inches(5.3), Inches(11), Inches(1.5), [
    {'text': '建议截图内容：', 'size': 13, 'color': GRAY_TEXT, 'bold': True, 'align': PP_ALIGN.CENTER},
    {'text': ' 教案生成页面（含模板选择）  ② 学生画像雷达图 + 时间轴  ③ 拍照答疑分步解析', 'size': 12, 'color': GRAY_TEXT, 'align': PP_ALIGN.CENTER, 'space_before': 10},
    {'text': '④ 错题本列表   家长端学情报告  ⑥ 离线模式演示', 'size': 12, 'color': GRAY_TEXT, 'align': PP_ALIGN.CENTER, 'space_before': 6},
])


# ============================================================
# 第6页：用户故事
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_top_bar(slide, '用户故事')
add_page_number(slide, 6)

stories = [
    ('👩‍🏫', '李老师', '35岁，云南山区村小数学老师',
     '教两个班87个学生，每天备课到深夜。\n现在30秒生成教案，5分钟编辑完成。\n备课时间从2小时缩短到30分钟。',
     ORANGE_PRIMARY, RGBColor(0xFF, 0xF5, 0xF0)),
    ('👦', '王小刚', '10岁，四年级，父母在外打工',
     '数学基础差，上课不敢问问题。\n用拍照答疑 + 错题本，\n期中考试从45分提高到68分。',
     GREEN, RGBColor(0xF0, 0xFD, 0xFB)),
    ('', '小刚奶奶', '62岁，只上过小学三年级',
     '不识字但能看懂报告：\n绿色箭头 = 进步，星星 = 表现好。\n每周和孙子一起看报告，心里踏实了。',
     PURPLE, RGBColor(0xF8, 0xF6, 0xFF)),
]

for i, (icon, name, role, story, color, bg) in enumerate(stories):
    x = Inches(0.8) + Inches(i * 4.1)
    card = add_rounded_rect(slide, x, Inches(1.2), Inches(3.8), Inches(5.8), bg)
    card.line.color.rgb = color
    card.line.width = Pt(1.5)

    add_text_box(slide, x + Inches(0.3), Inches(1.4), Inches(0.6), Inches(0.6),
                 icon, font_size=36)
    add_text_box(slide, x + Inches(1.0), Inches(1.4), Inches(2.5), Inches(0.4),
                 name, font_size=20, bold=True, color=color)
    add_text_box(slide, x + Inches(0.3), Inches(2.0), Inches(3.2), Inches(0.4),
                 role, font_size=11, color=GRAY_TEXT)

    add_shape_bg(slide, x + Inches(0.3), Inches(2.5), Inches(3.2), Inches(0.03), color)

    # 故事文本
    lines = story.split('\n')
    for j, line in enumerate(lines):
        add_text_box(slide, x + Inches(0.3), Inches(2.7) + Inches(j * 0.65), Inches(3.2), Inches(0.6),
                     line, font_size=13, color=DARK_TEXT)


# ============================================================
# 第7页：乡村场景适配
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_top_bar(slide, '乡村场景适配')
add_page_number(slide, 7)

adaptations = [
    ('📡', '离线能力', 'Service Worker 缓存\n教室断网也能用\n已生成内容完全离线可查', ORANGE_PRIMARY),
    ('📶', '低带宽优化', '纯前端架构\n无需下载大型安装包\n页面加载轻量快速', GREEN),
    ('💻', '零部署成本', '打开浏览器就能用\n不需要安装任何软件\n一台普通电脑即可运行', PURPLE),
    ('🔒', '数据安全', '所有数据存储在本地\n不上传云端\n保护学生隐私', RGBColor(0xFF, 0x8C, 0x42)),
    ('🔍', '大字体模式', '适配乡村老年教师\n清晰易读\n降低使用门槛', RGBColor(0x4E, 0xCD, 0xC4)),
]

for i, (icon, title, desc, color) in enumerate(adaptations):
    col = i % 3
    row = i // 3
    x = Inches(0.8) + Inches(col * 4.1)
    y = Inches(1.3) + Inches(row * 3.0)

    card = add_rounded_rect(slide, x, y, Inches(3.8), Inches(2.7), WHITE)
    card.line.color.rgb = color
    card.line.width = Pt(2)

    add_shape_bg(slide, x + Inches(0.05), y + Inches(0.05), Inches(3.7), Inches(0.06), color)

    add_text_box(slide, x + Inches(0.3), y + Inches(0.3), Inches(0.6), Inches(0.6),
                 icon, font_size=32)
    add_text_box(slide, x + Inches(1.0), y + Inches(0.35), Inches(2.5), Inches(0.4),
                 title, font_size=18, bold=True, color=color)

    add_shape_bg(slide, x + Inches(0.3), y + Inches(1.0), Inches(3.2), Inches(0.03), color)

    for j, line in enumerate(desc.split('\n')):
        add_text_box(slide, x + Inches(0.3), y + Inches(1.2) + Inches(j * 0.45), Inches(3.2), Inches(0.4),
                     f'• {line}', font_size=13, color=DARK_TEXT)


# ============================================================
# 第8页：效果对比
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_top_bar(slide, '效果对比')
add_page_number(slide, 8)

# 表格
headers = ['维度', '使用前', '使用后', '提升']
rows = [
    ['备课时间', '2小时/课', '30分钟/课', '⬇ 75%'],
    ['学生问题响应', '次日', '即时', '⚡ 即时'],
    ['家长了解学情', '家长会(1-2次/学期)', '随时查看', '📱 随时'],
    ['错题整理', '手动抄写', '自动归类', ' 自动'],
    ['因材施教', '凭经验', '数据驱动', '📊 精准'],
]

col_widths = [Inches(2.5), Inches(3.5), Inches(3.5), Inches(2.5)]
col_starts = [Inches(0.8)]
for w in col_widths[:-1]:
    col_starts.append(col_starts[-1] + w)

# 表头
for j, (header, cw, cs) in enumerate(zip(headers, col_widths, col_starts)):
    cell = add_rounded_rect(slide, cs, Inches(1.3), cw - Inches(0.05), Inches(0.6), ORANGE_PRIMARY)
    add_text_box(slide, cs + Inches(0.1), Inches(1.35), cw - Inches(0.2), Inches(0.5),
                 header, font_size=14, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

# 数据行
for i, row in enumerate(rows):
    bg = WHITE if i % 2 == 0 else ORANGE_BG
    for j, (cell_text, cw, cs) in enumerate(zip(row, col_widths, col_starts)):
        cell = add_rounded_rect(slide, cs, Inches(2.0) + Inches(i * 0.75), cw - Inches(0.05), Inches(0.65), bg)
        cell.line.fill.background()

        font_color = DARK_TEXT
        if j == 3:
            font_color = ORANGE_PRIMARY
        elif j == 2:
            font_color = GREEN

        add_text_box(slide, cs + Inches(0.1), Inches(2.05) + Inches(i * 0.75), cw - Inches(0.2), Inches(0.55),
                     cell_text, font_size=13, color=font_color,
                     bold=(j == 0), alignment=PP_ALIGN.CENTER)


# ============================================================
# 第9页：项目亮点总结
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_top_bar(slide, '项目亮点')
add_page_number(slide, 9)

highlights = [
    ('🎯', '真场景', '深入理解乡村教育痛点\n不是城市方案的简单移植\n从真实需求出发设计功能', ORANGE_PRIMARY),
    ('✅', '真可用', '纯前端 + 离线能力\n一台电脑就能跑\n零部署、零维护成本', GREEN),
    ('🔄', '真闭环', '教师→学生→家长\n三端数据互通\n完整的教育服务链路', PURPLE),
    ('❤️', '真温度', 'UI设计温暖克制\n不炫技，重实用\n每个功能都有人文关怀', RGBColor(0xFF, 0x8C, 0x42)),
]

for i, (icon, title, desc, color) in enumerate(highlights):
    x = Inches(0.8) + Inches(i * 3.1)
    card = add_rounded_rect(slide, x, Inches(1.3), Inches(2.8), Inches(5.5), WHITE)
    card.line.color.rgb = color
    card.line.width = Pt(2)

    # 顶部色块
    add_shape_bg(slide, x + Inches(0.05), Inches(1.35), Inches(2.7), Inches(1.5), color)

    add_text_box(slide, x + Inches(0.3), Inches(1.5), Inches(2.2), Inches(0.8),
                 icon, font_size=48, color=WHITE, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x + Inches(0.3), Inches(2.3), Inches(2.2), Inches(0.4),
                 title, font_size=22, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

    for j, line in enumerate(desc.split('\n')):
        add_text_box(slide, x + Inches(0.3), Inches(3.2) + Inches(j * 0.6), Inches(2.2), Inches(0.5),
                     f'• {line}', font_size=13, color=DARK_TEXT)


# ============================================================
# 第10页：未来展望
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_top_bar(slide, '未来展望')
add_page_number(slide, 10)

futures = [
    ('️', '方言语音答疑', '乡村学生普通话不标准\n支持方言识别，降低使用门槛', ORANGE_PRIMARY),
    ('', '教师社区', '教案共享，经验交流\n让优秀教学资源流动起来', GREEN),
    ('🧠', '自适应学习路径', '根据学生画像智能推荐\n个性化学习内容', PURPLE),
    ('🔗', '系统对接', '与学校现有系统对接\n成绩导入、课表同步', RGBColor(0xFF, 0x8C, 0x42)),
]

for i, (icon, title, desc, color) in enumerate(futures):
    col = i % 2
    row = i // 2
    x = Inches(0.8) + Inches(col * 6.2)
    y = Inches(1.3) + Inches(row * 2.8)

    card = add_rounded_rect(slide, x, y, Inches(5.8), Inches(2.5), WHITE)
    card.line.color.rgb = color
    card.line.width = Pt(2)

    add_text_box(slide, x + Inches(0.3), y + Inches(0.3), Inches(0.7), Inches(0.7),
                 icon, font_size=36)
    add_text_box(slide, x + Inches(1.1), y + Inches(0.35), Inches(4.5), Inches(0.5),
                 title, font_size=20, bold=True, color=color)

    add_shape_bg(slide, x + Inches(0.3), y + Inches(1.0), Inches(5.2), Inches(0.03), color)

    for j, line in enumerate(desc.split('\n')):
        add_text_box(slide, x + Inches(0.3), y + Inches(1.2) + Inches(j * 0.45), Inches(5.2), Inches(0.4),
                     f'• {line}', font_size=13, color=DARK_TEXT)


# ============================================================
# 第11页：致谢
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)

# 顶部装饰线
add_shape_bg(slide, 0, 0, W, Inches(0.08), ORANGE_PRIMARY)

# 主标题
add_text_box(slide, Inches(2), Inches(2.0), Inches(9), Inches(1.0),
             '感谢观看', font_size=48, bold=True, color=DARK_TEXT, alignment=PP_ALIGN.CENTER)

# 分隔线
add_shape_bg(slide, Inches(5.5), Inches(3.2), Inches(2.3), Inches(0.04), ORANGE_PRIMARY)

# 副标题
add_text_box(slide, Inches(2), Inches(3.5), Inches(9), Inches(0.6),
             '技术有温度 · 教育无距离', font_size=24, color=ORANGE_PRIMARY, alignment=PP_ALIGN.CENTER)

# 致谢内容
add_multi_text(slide, Inches(2), Inches(4.5), Inches(9), Inches(2.5), [
    {'text': '感谢阿里巴巴"小有可为"比赛提供平台', 'size': 16, 'color': GRAY_TEXT, 'align': PP_ALIGN.CENTER, 'space_before': 10},
    {'text': '感谢乡村教师的坚守与付出', 'size': 16, 'color': GRAY_TEXT, 'align': PP_ALIGN.CENTER, 'space_before': 10},
    {'text': '', 'size': 10},
    {'text': '乡村课堂AI助教', 'size': 20, 'color': ORANGE_PRIMARY, 'bold': True, 'align': PP_ALIGN.CENTER, 'space_before': 20},
    {'text': '让每一位乡村教师都有AI助手', 'size': 14, 'color': GRAY_TEXT, 'align': PP_ALIGN.CENTER, 'space_before': 8},
])

# 底部装饰线
add_shape_bg(slide, 0, H - Inches(0.08), W, Inches(0.08), ORANGE_PRIMARY)


# === 保存 ===
output_path = '/workspace/docs/乡村课堂AI助教_比赛PPT.pptx'
prs.save(output_path)
print(f'PPT 已生成: {output_path}')
print(f'共 {len(prs.slides)} 页')
