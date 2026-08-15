#!/usr/bin/env python3
"""生成乡村课堂AI助教项目说明PPT"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ── 颜色定义 ──
PRIMARY = RGBColor(0x1A, 0x73, 0xE8)       # 主蓝
PRIMARY_DARK = RGBColor(0x0D, 0x47, 0xA1)  # 深蓝
ACCENT = RGBColor(0x34, 0xA8, 0x53)        # 绿色
ACCENT2 = RGBColor(0xFB, 0xBC, 0x04)       # 黄色
ACCENT3 = RGBColor(0xEA, 0x43, 0x35)       # 红色
TEXT_DARK = RGBColor(0x20, 0x21, 0x24)     # 深灰文字
TEXT_LIGHT = RGBColor(0x5F, 0x63, 0x68)    # 浅灰文字
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BG_LIGHT = RGBColor(0xF8, 0xF9, 0xFA)      # 浅灰背景
BG_BLUE_LIGHT = RGBColor(0xE8, 0xF0, 0xFE) # 浅蓝背景

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
W = prs.slide_width
H = prs.slide_height

# ── 辅助函数 ──
def add_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_shape(slide, left, top, width, height, color, shape_type=MSO_SHAPE.RECTANGLE):
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_text_box(slide, left, top, width, height, text, font_size=18, color=TEXT_DARK, bold=False, alignment=PP_ALIGN.LEFT, font_name='Microsoft YaHei'):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_rounded_rect(slide, left, top, width, height, color, text, font_size=14, text_color=TEXT_DARK, bold=False):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(12)
    tf.margin_right = Pt(12)
    tf.margin_top = Pt(8)
    tf.margin_bottom = Pt(8)
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = text_color
    p.font.bold = bold
    p.font.name = 'Microsoft YaHei'
    p.alignment = PP_ALIGN.CENTER
    return shape

def add_bullet_list(slide, left, top, width, height, items, font_size=16, color=TEXT_DARK, spacing=Pt(6)):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = 'Microsoft YaHei'
        p.space_after = spacing
    return txBox

def add_top_bar(slide):
    """顶部蓝色装饰条"""
    add_shape(slide, 0, 0, W, Inches(0.08), PRIMARY)

def add_bottom_bar(slide, page_num, total=10):
    """底部信息栏"""
    bar = add_shape(slide, 0, H - Inches(0.5), W, Inches(0.5), PRIMARY)
    add_text_box(slide, Inches(0.5), H - Inches(0.45), Inches(5), Inches(0.4),
                 '乡村课堂AI助教  |  小有可为AI向善创新挑战赛', 12, WHITE)
    add_text_box(slide, W - Inches(2), H - Inches(0.45), Inches(1.5), Inches(0.4),
                 f'{page_num} / {total}', 12, WHITE, alignment=PP_ALIGN.RIGHT)

def add_section_title(slide, title, subtitle=''):
    """页面标题"""
    add_text_box(slide, Inches(0.8), Inches(0.3), Inches(10), Inches(0.7),
                 title, 32, PRIMARY_DARK, bold=True)
    if subtitle:
        add_text_box(slide, Inches(0.8), Inches(0.95), Inches(10), Inches(0.5),
                     subtitle, 16, TEXT_LIGHT)
    # 标题下划线
    add_shape(slide, Inches(0.8), Inches(1.5), Inches(2), Inches(0.06), PRIMARY)

# ═══════════════════════════════════════════
# 第1页：封面
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_bg(slide, WHITE)

# 左侧蓝色区域
add_shape(slide, 0, 0, Inches(6.5), H, PRIMARY)
# 装饰圆
circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(-1), Inches(-1), Inches(4), Inches(4))
circle.fill.solid()
circle.fill.fore_color.rgb = RGBColor(0x15, 0x65, 0xC0)
circle.line.fill.background()
circle2 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(4), Inches(5), Inches(3), Inches(3))
circle2.fill.solid()
circle2.fill.fore_color.rgb = RGBColor(0x15, 0x65, 0xC0)
circle2.line.fill.background()

# 封面标题
add_text_box(slide, Inches(1), Inches(2), Inches(5), Inches(1.2),
             '乡村课堂AI助教', 44, WHITE, bold=True)
add_text_box(slide, Inches(1), Inches(3.2), Inches(5), Inches(0.8),
             '让每一个乡村孩子都能享受AI教育', 20, RGBColor(0xBB, 0xDE, 0xFB))

# 右侧内容
add_text_box(slide, Inches(7.5), Inches(2), Inches(5), Inches(0.6),
             '小有可为 AI 向善创新挑战赛', 22, PRIMARY, bold=True)
add_text_box(slide, Inches(7.5), Inches(2.8), Inches(5), Inches(0.5),
             '参赛赛道：乡村教育 · 大山里的AI课', 16, TEXT_LIGHT)
add_text_box(slide, Inches(7.5), Inches(3.6), Inches(5), Inches(0.5),
             '作品类型：应用类（Web应用）', 16, TEXT_LIGHT)

# 底部标签
tags = ['纯前端架构', '三端协同', '多模态AI', '数据隔离', '低带宽友好']
for i, tag in enumerate(tags):
    add_rounded_rect(slide, Inches(7.5 + i * 1.1), Inches(4.5), Inches(1), Inches(0.45),
                     BG_BLUE_LIGHT, tag, 12, PRIMARY, bold=True)

add_top_bar(slide)
add_bottom_bar(slide, 1)

# ═══════════════════════════════════════════
# 第2页：痛点分析
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_top_bar(slide)
add_section_title(slide, '痛点分析', '乡村教育面临的真实困境')

pain_points = [
    ('📚', '教师备课负担重', '乡村教师往往身兼多科，\n缺乏AI工具辅助备课，\n教案质量参差不齐', ACCENT3),
    ('👨‍🎓', '学生个性化缺失', '大班教学难以因材施教，\n学生薄弱知识点难以\n精准定位和针对性训练', ACCENT2),
    ('👨‍👩‍', '家校沟通困难', '留守儿童多，家长外出务工，\n家校信息不对称，\n家长难以了解孩子学习情况', ACCENT),
    ('', '网络条件受限', '乡村网络带宽有限，\n传统云端AI应用\n难以流畅运行', PRIMARY),
]

for i, (icon, title, desc, color) in enumerate(pain_points):
    x = Inches(0.8 + i * 3.1)
    y = Inches(2)
    # 卡片背景
    card = add_shape(slide, x, y, Inches(2.8), Inches(4.5), BG_LIGHT)
    # 顶部色条
    add_shape(slide, x, y, Inches(2.8), Inches(0.08), color)
    # 图标
    add_text_box(slide, x + Inches(0.3), y + Inches(0.3), Inches(2.2), Inches(0.6),
                 icon, 36, color, alignment=PP_ALIGN.CENTER)
    # 标题
    add_text_box(slide, x + Inches(0.2), y + Inches(1.1), Inches(2.4), Inches(0.5),
                 title, 18, TEXT_DARK, bold=True, alignment=PP_ALIGN.CENTER)
    # 描述
    add_text_box(slide, x + Inches(0.2), y + Inches(1.8), Inches(2.4), Inches(2.5),
                 desc, 14, TEXT_LIGHT, alignment=PP_ALIGN.CENTER)

add_bottom_bar(slide, 2)

# ═══════════════════════════════════════════
# 第3页：产品简介
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_top_bar(slide)
add_section_title(slide, '产品简介', '面向乡村教育的AI助教平台，三端口协同')

# 三个端口卡片
ports = [
    ('👩‍', '教师端', '备课助手 · 学生画像 · 班级总览', PRIMARY),
    ('👨‍', '学生端', '拍照答疑 · 针对性练习 · 错题本', ACCENT),
    ('👨‍👩‍', '家长端', '学情报告 · 家校沟通', ACCENT2),
]

for i, (icon, title, desc, color) in enumerate(ports):
    x = Inches(0.8 + i * 4.1)
    y = Inches(2)
    card = add_shape(slide, x, y, Inches(3.8), Inches(3.5), BG_LIGHT)
    add_shape(slide, x, y, Inches(3.8), Inches(0.08), color)
    add_text_box(slide, x + Inches(0.3), y + Inches(0.3), Inches(3.2), Inches(0.7),
                 icon, 42, color, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x + Inches(0.3), y + Inches(1.2), Inches(3.2), Inches(0.5),
                 title, 24, TEXT_DARK, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x + Inches(0.3), y + Inches(1.9), Inches(3.2), Inches(0.5),
                 desc, 14, TEXT_LIGHT, alignment=PP_ALIGN.CENTER)

# 底部核心数据
data_items = [
    ('12,300+', '行代码'),
    ('7', '核心功能模块'),
    ('3', '端口架构'),
    ('100%', '纯前端实现'),
]
for i, (num, label) in enumerate(data_items):
    x = Inches(1 + i * 3.1)
    y = Inches(5.8)
    add_text_box(slide, x, y, Inches(2.5), Inches(0.5),
                 num, 28, PRIMARY, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x, y + Inches(0.5), Inches(2.5), Inches(0.4),
                 label, 14, TEXT_LIGHT, alignment=PP_ALIGN.CENTER)

add_bottom_bar(slide, 3)

# ═══════════════════════════════════════════
# 第4页：核心功能 - 教师端
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_top_bar(slide)
add_section_title(slide, '核心功能 · 教师端', 'AI赋能教学全流程')

features_teacher = [
    ('📝', 'AI备课助手', '智能生成教案，自动分层作业\n（基础/提高/拓展三层）', PRIMARY),
    ('📊', '学生画像', '五维雷达图展示学生能力\n归因分析 + 个性化建议', PRIMARY),
    ('📈', '班级总览', '成绩分布图、学习趋势图\n薄弱知识点分析', PRIMARY),
    ('👥', '学生管理', '批量导入/导出学生数据\n完整的学生信息管理', PRIMARY),
]

for i, (icon, title, desc, color) in enumerate(features_teacher):
    x = Inches(0.8 + i * 3.1)
    y = Inches(2)
    card = add_shape(slide, x, y, Inches(2.8), Inches(4.2), BG_BLUE_LIGHT)
    add_shape(slide, x + Inches(0.1), y + Inches(0.1), Inches(2.6), Inches(0.06), color)
    add_text_box(slide, x + Inches(0.2), y + Inches(0.3), Inches(2.4), Inches(0.6),
                 icon, 36, color, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x + Inches(0.2), y + Inches(1.1), Inches(2.4), Inches(0.5),
                 title, 18, TEXT_DARK, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x + Inches(0.2), y + Inches(1.8), Inches(2.4), Inches(2.2),
                 desc, 14, TEXT_LIGHT, alignment=PP_ALIGN.CENTER)

add_bottom_bar(slide, 4)

# ═══════════════════════════════════════════
# 第5页：核心功能 - 学生端
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_top_bar(slide)
add_section_title(slide, '核心功能 · 学生端', 'AI陪伴式学习体验')

features_student = [
    ('📷', '拍照答疑', '拍照/上传题目图片\nAI智能解析 + 分步讲解\n支持语音输入问题', ACCENT),
    ('🎯', '针对性练习', '基于学生画像推荐\n薄弱点专项训练\n难度自适应调整', ACCENT),
    ('📕', '错题本', '自动收录错题\n按知识点分类整理\n定期复习提醒', ACCENT),
]

for i, (icon, title, desc, color) in enumerate(features_student):
    x = Inches(0.8 + i * 4.1)
    y = Inches(2)
    card = add_shape(slide, x, y, Inches(3.8), Inches(4.2), RGBColor(0xE8, 0xF5, 0xE9))
    add_shape(slide, x + Inches(0.1), y + Inches(0.1), Inches(3.6), Inches(0.06), color)
    add_text_box(slide, x + Inches(0.3), y + Inches(0.3), Inches(3.2), Inches(0.7),
                 icon, 42, color, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x + Inches(0.3), y + Inches(1.2), Inches(3.2), Inches(0.5),
                 title, 22, TEXT_DARK, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x + Inches(0.3), y + Inches(1.9), Inches(3.2), Inches(2.2),
                 desc, 15, TEXT_LIGHT, alignment=PP_ALIGN.CENTER)

add_bottom_bar(slide, 5)

# ═══════════════════════════════════════════
# 第6页：核心功能 - 家长端
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_top_bar(slide)
add_section_title(slide, '核心功能 · 家长端', '让关爱跨越距离')

features_parent = [
    ('📋', '学情报告', '周报/月报自动生成\n可视化图表展示\n学习数据一目了然', ACCENT2),
    ('💬', '家校沟通', 'AI生成个性化沟通话术\n基于孩子真实学习数据\n降低沟通门槛', ACCENT2),
]

for i, (icon, title, desc, color) in enumerate(features_parent):
    x = Inches(0.8 + i * 6)
    y = Inches(2)
    card = add_shape(slide, x, y, Inches(5.5), Inches(4.2), RGBColor(0xFF, 0xF8, 0xE1))
    add_shape(slide, x + Inches(0.1), y + Inches(0.1), Inches(5.3), Inches(0.06), color)
    add_text_box(slide, x + Inches(0.5), y + Inches(0.3), Inches(4.5), Inches(0.7),
                 icon, 42, color, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x + Inches(0.5), y + Inches(1.2), Inches(4.5), Inches(0.5),
                 title, 22, TEXT_DARK, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x + Inches(0.5), y + Inches(1.9), Inches(4.5), Inches(2.2),
                 desc, 15, TEXT_LIGHT, alignment=PP_ALIGN.CENTER)

# 权限隔离说明
add_rounded_rect(slide, Inches(2), Inches(5.5), Inches(9.3), Inches(1.2), BG_LIGHT,
                 ' 权限隔离设计：家长端仅可查看已绑定孩子的学习数据，保障学生隐私安全', 16, PRIMARY_DARK, bold=True)

add_bottom_bar(slide, 6)

# ═══════════════════════════════════════════
# 第7页：技术亮点
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_top_bar(slide)
add_section_title(slide, '技术亮点', '创新技术方案')

highlights = [
    ('🤖', '多模态视觉AI', '拍照答疑采用通义千问VL Max\n视觉理解模型，支持图片OCR、\n数学公式识别、分步解题推理', PRIMARY),
    ('', '多API兼容', '支持阿里云通义千问、DeepSeek\n及任意OpenAI兼容API\n用户可自由切换和自定义模型', ACCENT),
    ('🔒', '三端权限隔离', '教师端看全量数据\n学生端仅看个人信息\n家长端仅看绑定孩子数据', ACCENT3),
    ('', '低带宽优化', '纯前端架构，无需后端服务器\n数据存储在浏览器本地\n适合乡村弱网络环境', ACCENT2),
]

for i, (icon, title, desc, color) in enumerate(highlights):
    x = Inches(0.8 + i * 3.1)
    y = Inches(2)
    card = add_shape(slide, x, y, Inches(2.8), Inches(4.5), BG_LIGHT)
    add_shape(slide, x, y, Inches(2.8), Inches(0.08), color)
    add_text_box(slide, x + Inches(0.2), y + Inches(0.3), Inches(2.4), Inches(0.6),
                 icon, 36, color, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x + Inches(0.2), y + Inches(1.1), Inches(2.4), Inches(0.5),
                 title, 17, TEXT_DARK, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x + Inches(0.2), y + Inches(1.8), Inches(2.4), Inches(2.5),
                 desc, 13, TEXT_LIGHT, alignment=PP_ALIGN.CENTER)

add_bottom_bar(slide, 7)

# ═══════════════════════════════════════════
# 第8页：技术方案
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_top_bar(slide)
add_section_title(slide, '技术方案', '纯前端架构，开箱即用')

# 架构图 - 三层
# 用户层
add_rounded_rect(slide, Inches(1), Inches(2), Inches(3.5), Inches(0.7), PRIMARY,
                 '👩‍🏫 教师端  |  👨‍🎓 学生端  |  👨‍👩‍👧 家长端', 16, WHITE, bold=True)

# 箭头
add_text_box(slide, Inches(5.5), Inches(2.1), Inches(2), Inches(0.5),
             '▼', 24, TEXT_LIGHT, alignment=PP_ALIGN.CENTER)

# 应用层
app_layer = add_shape(slide, Inches(1), Inches(2.9), Inches(11.3), Inches(1.2), BG_BLUE_LIGHT)
add_text_box(slide, Inches(1.2), Inches(2.95), Inches(11), Inches(0.4),
             '应用层', 14, PRIMARY, bold=True)
add_text_box(slide, Inches(1.2), Inches(3.35), Inches(11), Inches(0.6),
             '路由管理  |  数据存储  |  API服务  |  权限控制  |  组件库', 14, TEXT_DARK)

# 箭头
add_text_box(slide, Inches(5.5), Inches(4.2), Inches(2), Inches(0.5),
             '▼', 24, TEXT_LIGHT, alignment=PP_ALIGN.CENTER)

# 技术栈
tech_layer = add_shape(slide, Inches(1), Inches(4.8), Inches(11.3), Inches(1.8), BG_LIGHT)
add_text_box(slide, Inches(1.2), Inches(4.85), Inches(11), Inches(0.4),
             '技术栈', 14, PRIMARY, bold=True)

tech_items = [
    'HTML5 + CSS3 + Vanilla JS',
    'IndexedDB 本地数据库',
    'localStorage 配置存储',
    'OpenAI 兼容 API（多模型）',
    'Canvas 数据可视化',
    'Web Speech API 语音输入',
    'FileReader API 图片处理',
    '响应式设计（桌面+移动端）',
]
for i, item in enumerate(tech_items):
    col = i % 4
    row = i // 4
    x = Inches(1.3 + col * 2.8)
    y = Inches(5.25 + row * 0.55)
    add_text_box(slide, x, y, Inches(2.6), Inches(0.45),
                 f'● {item}', 12, TEXT_DARK)

add_bottom_bar(slide, 8)

# ═══════════════════════════════════════════
# 第9页：使用场景
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_top_bar(slide)
add_section_title(slide, '使用场景', '真实乡村课堂应用案例')

scenarios = [
    ('🌅', '场景一：备课', '乡村教师王老师要准备\n明天的数学课，打开备课助手，\n输入课题，AI自动生成\n完整教案和三层分层作业，\n节省2小时备课时间。', PRIMARY),
    ('📱', '场景二：答疑', '学生小明遇到一道几何题\n不会做，拍照上传，\nAI识别图形并分步讲解\n解题思路，就像有一个\n24小时在线的辅导老师。', ACCENT),
    ('💝', '场景三：家校沟通', '在外打工的李妈妈想了解\n孩子学习情况，打开家长端\n查看学情报告，AI还帮她\n生成与老师沟通的话术，\n让关爱不再缺席。', ACCENT2),
]

for i, (icon, title, desc, color) in enumerate(scenarios):
    x = Inches(0.8 + i * 4.1)
    y = Inches(2)
    card = add_shape(slide, x, y, Inches(3.8), Inches(4.5), BG_LIGHT)
    add_shape(slide, x, y, Inches(3.8), Inches(0.08), color)
    add_text_box(slide, x + Inches(0.3), y + Inches(0.3), Inches(3.2), Inches(0.7),
                 icon, 42, color, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x + Inches(0.3), y + Inches(1.2), Inches(3.2), Inches(0.5),
                 title, 20, TEXT_DARK, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x + Inches(0.3), y + Inches(1.9), Inches(3.2), Inches(2.5),
                 desc, 14, TEXT_LIGHT, alignment=PP_ALIGN.CENTER)

add_bottom_bar(slide, 9)

# ═══════════════════════════════════════════
# 第10页：项目成果与未来规划
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_top_bar(slide)
add_section_title(slide, '项目成果与未来规划')

# 左侧：项目成果
add_shape(slide, Inches(0.8), Inches(2), Inches(5.5), Inches(4.5), BG_BLUE_LIGHT)
add_text_box(slide, Inches(1), Inches(2.1), Inches(5), Inches(0.5),
             ' 项目成果', 20, PRIMARY_DARK, bold=True)

achievements = [
    '✅ 12,300+ 行代码，功能完整',
    '✅ 三端口架构：教师/学生/家长',
    '✅ 7大核心功能模块',
    '✅ 支持阿里云/DeepSeek/自定义API',
    '✅ 三端权限隔离，数据安全合规',
    '✅ 纯前端实现，开箱即用',
    '✅ 适配桌面端和移动端',
    '✅ 低带宽友好，适合乡村环境',
]
for i, item in enumerate(achievements):
    add_text_box(slide, Inches(1.2), Inches(2.7 + i * 0.45), Inches(5), Inches(0.4),
                 item, 14, TEXT_DARK)

# 右侧：未来规划
add_shape(slide, Inches(7), Inches(2), Inches(5.5), Inches(4.5), RGBColor(0xE8, 0xF5, 0xE9))
add_text_box(slide, Inches(7.2), Inches(2.1), Inches(5), Inches(0.5),
             '🚀 未来规划', 20, RGBColor(0x1B, 0x5E, 0x20), bold=True)

future = [
    '📌 接入更多学科题库资源',
    ' 支持离线模式，无网也能用',
    '📌 增加教师间教案共享社区',
    '📌 引入AI口语练习功能',
    '📌 对接学校教务系统',
    '📌 多语言支持（少数民族地区）',
    '📌 家长端增加学习提醒推送',
    '📌 开源社区共建',
]
for i, item in enumerate(future):
    add_text_box(slide, Inches(7.4), Inches(2.7 + i * 0.45), Inches(5), Inches(0.4),
                 item, 14, TEXT_DARK)

# 底部标语
add_text_box(slide, Inches(2), Inches(6.6), Inches(9.3), Inches(0.6),
             '用AI点亮乡村教育，让每个孩子都能看见未来', 22, PRIMARY, bold=True, alignment=PP_ALIGN.CENTER)

add_bottom_bar(slide, 10)

# ── 保存 ──
output_path = '/workspace/乡村课堂AI助教-项目说明.pptx'
prs.save(output_path)
print(f'PPT 已生成：{output_path}')
print(f'文件大小：{os.path.getsize(output_path) / 1024:.1f} KB')
print(f'共 {len(prs.slides)} 页')
